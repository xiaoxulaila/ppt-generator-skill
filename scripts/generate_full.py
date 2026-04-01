#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
内容丰满版 PPT 生成器
AI发展趋势报告
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# 配色方案
COLORS = {
    'primary': RGBColor(30, 39, 97),      # 深蓝
    'secondary': RGBColor(202, 220, 252),  # 冰蓝
    'accent': RGBColor(88, 166, 255),     # 亮蓝
    'white': RGBColor(255, 255, 255),
    'text': RGBColor(50, 50, 80),
    'gray': RGBColor(120, 120, 140),
}

def add_slide_background(slide, prs, color=None):
    """添加背景"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color or COLORS['white']
    bg.line.fill.background()
    return bg

def add_decorator_circles(slide, prs):
    """添加装饰圆形"""
    # 右上角
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1), Inches(4), Inches(4))
    c1.fill.solid()
    c1.fill.fore_color.rgb = COLORS['secondary']
    c1.line.fill.background()
    # 左下角
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(5.5), Inches(3), Inches(3))
    c2.fill.solid()
    c2.fill.fore_color.rgb = COLORS['secondary']
    c2.line.fill.background()

def add_text(slide, text, x, y, w, h, size=16, bold=False, color=None, align='left'):
    """添加文字"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = 'Microsoft YaHei'
    p.font.color.rgb = color or COLORS['text']
    if align == 'center':
        p.alignment = PP_ALIGN.CENTER
    return box

def add_content_point(slide, icon, text, x, y, color):
    """添加内容要点"""
    # 圆点
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y + 0.1), Inches(0.3), Inches(0.3))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    # 文字
    add_text(slide, text, x + 0.5, y, 10, 0.8, size=14, color=COLORS['text'])

def create_full_ppt(output_path=None):
    """生成内容丰满的PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ========== 第1页：封面 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs, COLORS['primary'])
    add_decorator_circles(slide, prs)
    
    add_text(slide, "AI发展趋势报告", 1, 2.5, 11, 1.5, size=48, bold=True, color=COLORS['white'], align='center')
    add_text(slide, "2025-2026年行业深度洞察", 1, 4.2, 11, 0.8, size=22, color=COLORS['secondary'], align='center')
    add_text(slide, "基于 Gartner、IDC、艾瑞咨询等权威报告", 1, 5.5, 11, 0.5, size=14, color=COLORS['secondary'], align='center')
    
    # ========== 第2页：目录 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs)
    add_decorator_circles(slide, prs)
    
    add_text(slide, "目录 CONTENTS", 0.8, 0.4, 5, 0.8, size=32, bold=True, color=COLORS['primary'])
    
    toc_items = [
        ("01", "全球AI市场规模", "2025年达4500亿美元，年增长20%"),
        ("02", "技术演进趋势", "大模型多模态、端侧AI、AI Agent成为新风口"),
        ("03", "应用场景爆发", "医疗、教育、制造、金融等领域全面渗透"),
        ("04", "行业案例分析", "国内外头部企业AI布局详解"),
        ("05", "未来展望与建议", "2026年趋势预测与行动建议"),
    ]
    
    for i, (num, title, desc) in enumerate(toc_items):
        y = 1.5 + i * 1.1
        # 序号
        add_text(slide, num, 0.8, y, 0.8, 0.6, size=24, bold=True, color=COLORS['accent'])
        # 标题
        add_text(slide, title, 1.8, y, 4, 0.5, size=18, bold=True, color=COLORS['text'])
        # 描述
        add_text(slide, desc, 1.8, y + 0.4, 9, 0.5, size=12, color=COLORS['gray'])
    
    # ========== 第3页：全球市场规模 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs)
    add_decorator_circles(slide, prs)
    
    add_text(slide, "全球AI市场规模", 0.8, 0.4, 6, 0.8, size=32, bold=True, color=COLORS['primary'])
    add_text(slide, "Global AI Market Size", 0.8, 1, 6, 0.5, size=14, color=COLORS['gray'])
    
    # 数据卡片
    data_items = [
        ("4500亿", "美元", "2025年全球AI市场总规模", COLORS['primary']),
        ("20%", "CAGR", "年复合增长率", COLORS['accent']),
        ("67%", "中国企业", "已在业务中采用AI技术", COLORS['secondary']),
    ]
    
    for i, (value, unit, desc, color) in enumerate(data_items):
        x = 0.8 + i * 4.2
        # 卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2), Inches(3.8), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()
        # 数值
        add_text(slide, value, x + 0.2, 2.3, 3.4, 0.9, size=36, bold=True, color=COLORS['white'], align='center')
        add_text(slide, unit, x + 0.2, 3.1, 3.4, 0.5, size=16, color=COLORS['white'], align='center')
        add_text(slide, desc, x + 0.2, 3.5, 3.4, 0.5, size=11, color=COLORS['secondary'], align='center')
    
    # 详细说明
    details = [
        "• 生成式AI市场增长最为迅猛，预计2025年占比超30%",
        "• 中国AI市场增速领先全球，2025年规模预计达1500亿人民币",
        "• 企业级AI解决方案市场年增长率达35%，成为增长主力",
        "• AI芯片市场规模预计2025年突破1000亿美元",
    ]
    for i, detail in enumerate(details):
        add_text(slide, detail, 0.8, 4.5 + i * 0.55, 11, 0.5, size=13, color=COLORS['text'])
    
    # ========== 第4页：技术演进趋势 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs)
    add_decorator_circles(slide, prs)
    
    add_text(slide, "技术演进趋势", 0.8, 0.4, 6, 0.8, size=32, bold=True, color=COLORS['primary'])
    add_text(slide, "Technology Evolution Trends", 0.8, 1, 6, 0.5, size=14, color=COLORS['gray'])
    
    tech_items = [
        (COLORS['primary'], "大模型多模态", "GPT-5、Gemini 2.0等旗舰模型全面支持文本、图像、视频、音频多模态理解与生成，实现真正的跨模态智能。"),
        (COLORS['accent'], "端侧AI加速普及", "搭载NPU的智能手机、PC设备快速增长。到2026年，预计60%以上的新出货设备将具备本地AI处理能力。"),
        (COLORS['secondary'], "AI Agent成为新风口", "AutoGPT、BabyAGI等Agent项目爆发式增长。AI从\"工具\"进化为\"助手\"，能够自主规划、执行复杂任务。"),
        (COLORS['gray'], "开源生态崛起", "Llama 3、Mistral等开源模型性能逼近GPT-4，推动AI技术民主化，降低企业应用门槛。"),
    ]
    
    for i, (color, title, desc) in enumerate(tech_items):
        y = 1.8 + i * 1.3
        # 左侧色块
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(0.15), Inches(1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        # 标题
        add_text(slide, title, 1.2, y, 4, 0.5, size=16, bold=True, color=COLORS['text'])
        # 描述
        add_text(slide, desc, 1.2, y + 0.45, 11, 0.7, size=12, color=COLORS['gray'])
    
    # ========== 第5页：应用场景爆发 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs)
    add_decorator_circles(slide, prs)
    
    add_text(slide, "应用场景爆发", 0.8, 0.4, 6, 0.8, size=32, bold=True, color=COLORS['primary'])
    add_text(slide, "Application Scenarios Explosion", 0.8, 1, 6, 0.5, size=14, color=COLORS['gray'])
    
    # 四个应用领域
    apps = [
        (COLORS['primary'], "🏥 医疗健康", "AI辅助诊断准确率达95%\n智能药物研发周期缩短60%\n2025年AI医疗市场达150亿美元"),
        (COLORS['accent'], "📚 教育培训", "个性化学习效率提升60%\nAI tutoring覆盖1亿+学生\n作业批改时间减少70%"),
        (COLORS['secondary'], "🏭 智能制造", "智能质检效率提升40%\n预测性维护降低停机30%\n数字孪生应用渗透率达45%"),
        (COLORS['gray'], "💰 金融服务", "智能风控坏账率降低30%\nAI客服替代率达80%\n智能投顾管理规模破万亿"),
    ]
    
    for i, (color, title, desc) in enumerate(apps):
        row = i // 2
        col = i % 2
        x = 0.8 + col * 6.2
        y = 1.6 + row * 2.7
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.8), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()
        
        add_text(slide, title, x + 0.3, y + 0.3, 5, 0.6, size=18, bold=True, color=COLORS['white'])
        add_text(slide, desc, x + 0.3, y + 1, 5.2, 1.2, size=12, color=COLORS['white'])
    
    # ========== 第6页：行业案例分析 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs)
    add_decorator_circles(slide, prs)
    
    add_text(slide, "行业案例分析", 0.8, 0.4, 6, 0.8, size=32, bold=True, color=COLORS['primary'])
    add_text(slide, "Industry Case Studies", 0.8, 1, 6, 0.5, size=14, color=COLORS['gray'])
    
    cases = [
        ("OpenAI", "ChatGPT月活超2亿，企业客户超100万家，GPT-4 API调用量年增长10倍。"),
        ("Google", "Gemini 1.5 Pro上下文窗口达100万token，Bard累计服务超1亿用户。"),
        ("Microsoft", "Copilot全面集成Office 365，用户超150万，企业效率提升超40%。"),
        ("字节跳动", "AI推荐算法驱动抖音增长，AI内容生成覆盖80%短视频创作。"),
        ("华为", "盘古大模型赋能多个行业，AI算力基础设施全球前三。"),
        ("阿里巴巴", "通义千问服务企业超10万家，AI电商GMV贡献超15%。"),
    ]
    
    for i, (company, desc) in enumerate(cases):
        row = i // 2
        col = i % 2
        x = 0.8 + col * 6.2
        y = 1.6 + row * 1.8
        
        # 公司名
        add_text(slide, company, x, y, 5, 0.5, size=16, bold=True, color=COLORS['primary'])
        # 描述
        add_text(slide, desc, x, y + 0.5, 5.8, 1, size=12, color=COLORS['text'])
        # 分隔线
        if col == 0:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 5.9), Inches(y), Inches(0.05), Inches(1.5))
            line.fill.solid()
            line.fill.fore_color.rgb = COLORS['secondary']
            line.line.fill.background()
    
    # ========== 第7页：未来展望 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs, COLORS['primary'])
    add_decorator_circles(slide, prs)
    
    add_text(slide, "2026年趋势预测", 0.8, 0.4, 6, 0.8, size=32, bold=True, color=COLORS['white'])
    add_text(slide, "Future Outlook 2026", 0.8, 1, 6, 0.5, size=14, color=COLORS['secondary'])
    
    outlooks = [
        "AI Agent将成为个人和企业的标配助手",
        "多模态AI应用全面爆发，视频、3D生成常态化",
        "端侧AI普及，隐私保护成为核心竞争力",
        "AI原生应用崛起，颠覆传统软件交互方式",
        "AI监管框架逐步完善，合规成为必备能力",
    ]
    
    for i, text in enumerate(outlooks):
        y = 1.8 + i * 1
        # 序号
        add_text(slide, f"0{i+1}", 0.8, y, 0.6, 0.6, size=20, bold=True, color=COLORS['accent'])
        # 内容
        add_text(slide, text, 1.6, y, 10, 0.6, size=18, color=COLORS['white'])
    
    # ========== 第8页：总结 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs, COLORS['primary'])
    add_decorator_circles(slide, prs)
    
    add_text(slide, "总结与建议", 0.8, 0.4, 6, 0.8, size=32, bold=True, color=COLORS['white'])
    
    suggestions = [
        ("立即行动", "不要等待，先从小场景切入AI应用"),
        ("选对场景", "优先选择ROI明确、痛点清晰的场景"),
        ("数据为王", "高质量数据是AI落地成功的关键"),
        ("人才培养", "组建AI团队或与专业机构合作"),
        ("持续迭代", "AI技术演进快，保持学习和迭代"),
    ]
    
    for i, (title, desc) in enumerate(suggestions):
        y = 1.5 + i * 1.1
        # 标题
        add_text(slide, title, 0.8, y, 2, 0.5, size=16, bold=True, color=COLORS['accent'])
        # 描述
        add_text(slide, desc, 2.8, y, 9, 0.5, size=15, color=COLORS['white'])
    
    # ========== 第9页：结束页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs, COLORS['primary'])
    
    add_text(slide, "感谢观看", 0, 2.8, 13.333, 1.5, size=52, bold=True, color=COLORS['white'], align='center')
    add_text(slide, "THANK YOU", 0, 4.2, 13.333, 0.8, size=20, color=COLORS['secondary'], align='center')
    add_text(slide, "关注我们，获取更多AI行业洞察", 0, 5.5, 13.333, 0.5, size=14, color=COLORS['secondary'], align='center')
    
    # 保存
    if output_path is None:
        output_path = os.path.join(os.path.expanduser("~"), "Desktop", "AI发展趋势报告_丰满版.pptx")
    
    prs.save(output_path)
    print(f"✅ 内容丰满版PPT已生成: {output_path}")
    return output_path

if __name__ == "__main__":
    output = create_full_ppt()
    print(f"\n📁 文件位置: {output}")
