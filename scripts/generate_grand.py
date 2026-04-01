#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
大气高端风格 PPT 生成器
AI发展趋势报告
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# 高端大气配色
COLORS = {
    'bg_dark': RGBColor(10, 15, 30),     # 深邃黑蓝
    'bg_gradient': RGBColor(20, 30, 60), # 渐变深蓝
    'accent_gold': RGBColor(212, 175, 55), # 金色
    'accent_blue': RGBColor(70, 130, 200), # 宝石蓝
    'white': RGBColor(255, 255, 255),
    'light_blue': RGBColor(180, 200, 240),
    'gray': RGBColor(150, 160, 180),
    'dark_gray': RGBColor(80, 90, 110),
}

def add_slide_background(slide, prs, dark=True):
    """添加背景"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    if dark:
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLORS['bg_dark']
    else:
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()
    return bg

def add_gradient_bar(slide, prs, y=0, height=0.1):
    """添加渐变装饰条"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(y), prs.slide_width, Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['accent_gold']
    bar.line.fill.background()

def add_decorator_lines(slide, prs):
    """添加装饰线条"""
    # 底部金色线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.8), Inches(2), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent_gold']
    line.line.fill.background()

def add_text(slide, text, x, y, w, h, size=16, bold=False, color=None, align='left'):
    """添加文字"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = 'Microsoft YaHei'
    p.font.color.rgb = color or COLORS['white']
    if align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT
    return box

def create_grand_ppt(output_path=None):
    """生成大气高端PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ========== 第1页：封面 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs, dark=True)
    
    # 顶部金色装饰线
    add_gradient_bar(slide, prs, 0, 0.05)
    
    # 底部金色装饰线
    add_gradient_bar(slide, prs, 7.4, 0.05)
    
    # 左侧装饰条
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2), Inches(0.3), Inches(3.5))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = COLORS['accent_gold']
    left_bar.line.fill.background()
    
    # 主标题
    add_text(slide, "AI", 1, 2.2, 3, 1.5, size=72, bold=True, color=COLORS['white'], align='left')
    add_text(slide, "发展趋势报告", 3.5, 3.2, 7, 1, size=40, bold=True, color=COLORS['accent_gold'], align='left')
    
    # 副标题
    add_text(slide, "2025-2026年行业深度洞察", 1, 4.5, 10, 0.6, size=20, color=COLORS['light_blue'], align='left')
    
    # 来源
    add_text(slide, "基于 Gartner | IDC | 麦肯锡 | 艾瑞咨询", 1, 5.5, 10, 0.5, size=14, color=COLORS['gray'], align='left')
    
    # 右下角装饰
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4), Inches(4), Inches(4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS['bg_gradient']
    circle.line.fill.background()
    
    # ========== 第2页：目录 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs, dark=True)
    add_gradient_bar(slide, prs, 0, 0.05)
    add_gradient_bar(slide, prs, 7.4, 0.05)
    
    add_text(slide, "CONTENTS", 0.8, 0.5, 5, 0.8, size=14, color=COLORS['gray'])
    add_text(slide, "目录", 0.8, 1, 5, 1, size=40, bold=True, color=COLORS['white'])
    
    toc_items = [
        ("01", "全球AI市场规模", "市场规模、增长率、区域分布"),
        ("02", "技术演进趋势", "大模型、端侧AI、AI Agent"),
        ("03", "应用场景分析", "医疗、教育、制造、金融"),
        ("04", "行业案例研究", "头部企业AI战略布局"),
        ("05", "未来趋势展望", "2026年预测与建议"),
    ]
    
    for i, (num, title, desc) in enumerate(toc_items):
        y = 2.3 + i * 0.95
        
        # 序号
        add_text(slide, num, 0.8, y, 1, 0.6, size=28, bold=True, color=COLORS['accent_gold'])
        
        # 标题
        add_text(slide, title, 2, y + 0.1, 5, 0.5, size=20, bold=True, color=COLORS['white'])
        
        # 描述
        add_text(slide, desc, 2, y + 0.5, 5, 0.4, size=12, color=COLORS['gray'])
    
    # 右侧装饰
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9), Inches(1), Inches(4), Inches(6))
    rect.fill.solid()
    rect.fill.fore_color.rgb = COLORS['bg_gradient']
    rect.line.fill.background()
    
    # ========== 第3页：全球市场规模 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs, dark=True)
    add_gradient_bar(slide, prs, 0, 0.05)
    
    add_text(slide, "01", 0.8, 0.3, 1, 0.6, size=16, color=COLORS['accent_gold'])
    add_text(slide, "全球AI市场规模", 0.8, 0.7, 6, 0.8, size=36, bold=True, color=COLORS['white'])
    add_text(slide, "Global AI Market Size", 0.8, 1.4, 6, 0.4, size=14, color=COLORS['gray'])
    
    # 三个大数据卡片
    data_cards = [
        ("4500", "亿美元", "2025年全球AI市场总规模", COLORS['accent_gold']),
        ("20%", "CAGR", "年复合增长率预测", COLORS['accent_blue']),
        ("67%", "企业", "已在业务中采用AI技术", COLORS['light_blue']),
    ]
    
    for i, (value, unit, desc, color) in enumerate(data_cards):
        x = 0.8 + i * 4.2
        
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.2), Inches(3.8), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['bg_gradient']
        card.line.color.rgb = color
        card.line.width = Pt(2)
        
        # 数值
        add_text(slide, value, x + 0.2, 2.5, 3.4, 1, size=48, bold=True, color=color, align='center')
        # 单位
        add_text(slide, unit, x + 0.2, 3.4, 3.4, 0.5, size=18, color=COLORS['white'], align='center')
        # 描述
        add_text(slide, desc, x + 0.2, 4, 3.4, 0.5, size=12, color=COLORS['gray'], align='center')
    
    # 详细数据
    details = [
        "• 生成式AI市场增长迅猛，预计2025年占比超30%，成为最大细分市场",
        "• 中国AI市场增速领先全球，2025年规模预计达1500亿人民币",
        "• 企业级AI解决方案市场年增长率达35%，2028年市场规模将突破1万亿美元",
        "• AI芯片市场规模高速增长，2025年预计突破1000亿美元",
    ]
    
    for i, detail in enumerate(details):
        y = 5 + i * 0.5
        add_text(slide, detail, 0.8, y, 12, 0.5, size=13, color=COLORS['light_blue'])
    
    add_gradient_bar(slide, prs, 7.4, 0.05)
    
    # ========== 第4页：技术演进趋势 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs, dark=True)
    add_gradient_bar(slide, prs, 0, 0.05)
    
    add_text(slide, "02", 0.8, 0.3, 1, 0.6, size=16, color=COLORS['accent_gold'])
    add_text(slide, "技术演进趋势", 0.8, 0.7, 6, 0.8, size=36, bold=True, color=COLORS['white'])
    add_text(slide, "Technology Evolution Trends", 0.8, 1.4, 6, 0.4, size=14, color=COLORS['gray'])
    
    tech_items = [
        ("大模型多模态", "GPT-5、Gemini 2.0等旗舰模型全面支持文本、图像、视频、音频多模态理解与生成，实现真正的跨模态智能。"),
        ("端侧AI普及", "搭载NPU的智能设备快速增长。2026年预计60%以上新出货设备将具备本地AI处理能力。"),
        ("AI Agent爆发", "AutoGPT、BabyAGI等Agent项目爆发式增长。AI从\"工具\"进化为\"助手\"，自主规划和执行复杂任务。"),
        ("开源生态崛起", "Llama 3、Mistral等开源模型性能逼近GPT-4，推动AI技术民主化，降低企业应用门槛。"),
    ]
    
    for i, (title, desc) in enumerate(tech_items):
        y = 2.2 + i * 1.25
        
        # 左侧金色标记
        mark = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(0.08), Inches(0.9))
        mark.fill.solid()
        mark.fill.fore_color.rgb = COLORS['accent_gold']
        mark.line.fill.background()
        
        # 标题
        add_text(slide, title, 1.1, y, 4, 0.5, size=18, bold=True, color=COLORS['white'])
        # 描述
        add_text(slide, desc, 1.1, y + 0.45, 11, 0.7, size=12, color=COLORS['gray'])
    
    add_gradient_bar(slide, prs, 7.4, 0.05)
    
    # ========== 第5页：应用场景 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs, dark=True)
    add_gradient_bar(slide, prs, 0, 0.05)
    
    add_text(slide, "03", 0.8, 0.3, 1, 0.6, size=16, color=COLORS['accent_gold'])
    add_text(slide, "应用场景分析", 0.8, 0.7, 6, 0.8, size=36, bold=True, color=COLORS['white'])
    add_text(slide, "Application Scenarios", 0.8, 1.4, 6, 0.4, size=14, color=COLORS['gray'])
    
    apps = [
        (COLORS['accent_gold'], "医疗健康", "AI辅助诊断准确率达95%\n智能药物研发周期缩短60%\n2025年AI医疗市场达150亿美元"),
        (COLORS['accent_blue'], "教育培训", "个性化学习效率提升60%\nAI tutoring覆盖1亿+学生\n作业批改时间减少70%"),
        (COLORS['light_blue'], "智能制造", "智能质检效率提升40%\n预测性维护降低停机30%\n数字孪生应用渗透率达45%"),
        (COLORS['gray'], "金融服务", "智能风控坏账率降低30%\nAI客服替代率达80%\n智能投顾管理规模破万亿"),
    ]
    
    for i, (color, title, desc) in enumerate(apps):
        row = i // 2
        col = i % 2
        x = 0.8 + col * 6.2
        y = 2.2 + row * 2.4
        
        # 卡片
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(5.8), Inches(2.1))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['bg_gradient']
        card.line.color.rgb = color
        card.line.width = Pt(2)
        
        # 标题
        add_text(slide, title, x + 0.3, y + 0.3, 5, 0.5, size=20, bold=True, color=color)
        # 描述
        add_text(slide, desc, x + 0.3, y + 0.9, 5.2, 1, size=13, color=COLORS['light_blue'])
    
    add_gradient_bar(slide, prs, 7.4, 0.05)
    
    # ========== 第6页：行业案例 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs, dark=True)
    add_gradient_bar(slide, prs, 0, 0.05)
    
    add_text(slide, "04", 0.8, 0.3, 1, 0.6, size=16, color=COLORS['accent_gold'])
    add_text(slide, "行业案例研究", 0.8, 0.7, 6, 0.8, size=36, bold=True, color=COLORS['white'])
    add_text(slide, "Industry Case Studies", 0.8, 1.4, 6, 0.4, size=14, color=COLORS['gray'])
    
    cases = [
        ("OpenAI", "ChatGPT月活超2亿，企业客户超100万家，GPT-4 API调用量年增长10倍。"),
        ("Google", "Gemini 1.5 Pro上下文窗口达100万token，Bard累计服务超1亿用户。"),
        ("Microsoft", "Copilot全面集成Office 365，用户超150万，企业效率提升超40%。"),
        ("字节跳动", "AI推荐算法驱动抖音增长，AI内容生成覆盖80%短视频创作。"),
        ("华为", "盘古大模型赋能多个行业，AI算力基础设施全球前三。"),
        ("阿里巴巴", "通义千问服务企业超10万家，AI电商GMV贡献超15%。"),
    ]
    
    for i, (company, desc) in enumerate(cases):
        row = i // 2
        col = i % 2
        x = 0.8 + col * 6.2
        y = 2 + row * 1.7
        
        # 公司名
        add_text(slide, company, x, y, 5, 0.5, size=20, bold=True, color=COLORS['accent_gold'])
        # 描述
        add_text(slide, desc, x, y + 0.6, 5.8, 0.8, size=13, color=COLORS['light_blue'])
    
    add_gradient_bar(slide, prs, 7.4, 0.05)
    
    # ========== 第7页：未来展望 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs, dark=True)
    add_gradient_bar(slide, prs, 0, 0.05)
    
    add_text(slide, "05", 0.8, 0.3, 1, 0.6, size=16, color=COLORS['accent_gold'])
    add_text(slide, "未来趋势展望", 0.8, 0.7, 6, 0.8, size=36, bold=True, color=COLORS['white'])
    add_text(slide, "Future Outlook 2026", 0.8, 1.4, 6, 0.4, size=14, color=COLORS['gray'])
    
    outlooks = [
        "AI Agent将成为个人和企业的标配助手",
        "多模态AI应用全面爆发，视频、3D生成常态化",
        "端侧AI普及，隐私保护成为核心竞争力",
        "AI原生应用崛起，颠覆传统软件交互方式",
        "AI监管框架逐步完善，合规成为必备能力",
    ]
    
    for i, text in enumerate(outlooks):
        y = 2.2 + i * 0.95
        
        # 序号圆圈
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y + 0.1), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['accent_gold']
        circle.line.fill.background()
        
        add_text(slide, str(i + 1), 0.8, y + 0.15, 0.5, 0.4, size=16, bold=True, color=COLORS['bg_dark'], align='center')
        add_text(slide, text, 1.5, y + 0.15, 10, 0.5, size=18, color=COLORS['white'])
    
    add_gradient_bar(slide, prs, 7.4, 0.05)
    
    # ========== 第8页：总结 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs, dark=True)
    add_gradient_bar(slide, prs, 0, 0.05)
    
    add_text(slide, "总结与建议", 0.8, 0.5, 6, 1, size=40, bold=True, color=COLORS['white'])
    add_text(slide, "Summary & Recommendations", 0.8, 1.3, 6, 0.5, size=14, color=COLORS['gray'])
    
    suggestions = [
        ("立即行动", "不要等待，先从小场景切入AI应用，快速验证价值"),
        ("选对场景", "优先选择ROI明确、痛点清晰的场景，确保快速见效"),
        ("数据为王", "高质量数据是AI落地成功的关键，加大数据投入"),
        ("人才培养", "组建AI团队或与专业机构合作，建立AI能力"),
        ("持续迭代", "AI技术演进快，保持学习和迭代，建立技术储备"),
    ]
    
    for i, (title, desc) in enumerate(suggestions):
        y = 2.3 + i * 0.95
        
        # 标题
        add_text(slide, title, 0.8, y, 2.5, 0.5, size=18, bold=True, color=COLORS['accent_gold'])
        # 描述
        add_text(slide, desc, 3.5, y + 0.1, 9, 0.5, size=15, color=COLORS['light_blue'])
    
    add_gradient_bar(slide, prs, 7.4, 0.05)
    
    # ========== 第9页：结束页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs, dark=True)
    
    # 顶部金色线
    add_gradient_bar(slide, prs, 0, 0.08)
    # 底部金色线
    add_gradient_bar(slide, prs, 7.42, 0.08)
    
    # 左侧金色装饰
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), Inches(0.4), Inches(2.5))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = COLORS['accent_gold']
    left_bar.line.fill.background()
    
    add_text(slide, "THANK YOU", 1, 2.8, 12, 1, size=56, bold=True, color=COLORS['white'], align='center')
    add_text(slide, "感谢观看", 1, 4, 12, 0.8, size=28, color=COLORS['accent_gold'], align='center')
    add_text(slide, "关注我们，获取更多AI行业洞察", 1, 5, 12, 0.5, size=16, color=COLORS['gray'], align='center')
    
    # 右下角装饰
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4.5), Inches(4), Inches(4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS['bg_gradient']
    circle.line.fill.background()
    
    # 保存
    if output_path is None:
        output_path = os.path.join(os.path.expanduser("~"), "Desktop", "AI发展趋势报告_大气版.pptx")
    
    prs.save(output_path)
    print(f"✅ 大气高端版PPT已生成: {output_path}")
    return output_path

if __name__ == "__main__":
    output = create_grand_ppt()
    print(f"\n📁 文件位置: {output}")
