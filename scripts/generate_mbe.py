#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
MBE插画风格 PPT 生成器
大学新生开学第一课
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# MBE 风格配色
MBE_COLORS = {
    'bg': RGBColor(255, 255, 255),        # 白色背景
    'border': RGBColor(0, 0, 0),          # 黑色描边
    'text': RGBColor(26, 26, 26),         # 深灰文字
    'yellow': RGBColor(255, 214, 0),      # 明黄
    'purple': RGBColor(156, 39, 176),     # 亮紫
    'red': RGBColor(244, 67, 54),         # 正红
    'blue': RGBColor(33, 150, 243),       # 天蓝
    'green': RGBColor(76, 175, 80),       # 草绿
    'orange': RGBColor(255, 152, 0),       # 橙色
    'gray': RGBColor(128, 128, 128),       # 灰色
}

def add_decorative_elements(slide, prs, colors_list):
    """添加装饰元素"""
    # 右下角黄色圆形
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, 
        Inches(10.5), Inches(5.5), 
        Inches(3), Inches(3)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = MBE_COLORS['yellow']
    # 无描边
    circle.line.fill.background()
    
    # 左上角紫色圆形
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-0.5), Inches(-0.5),
        Inches(2), Inches(2)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = MBE_COLORS['purple']
    circle2.line.fill.background()
    
    # 中间装饰星星（用小圆形代替）
    for i, (x, y) in enumerate([(11, 1.5), (11.5, 2.5), (10.8, 3.5)]):
        star = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y),
            Inches(0.3), Inches(0.3)
        )
        star.fill.solid()
        star.fill.fore_color.rgb = colors_list[i % len(colors_list)]
        star.line.fill.background()

def create_mbe_ppt(output_path=None):
    """生成 MBE 风格 PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ===== 第1页：封面 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 白色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = MBE_COLORS['bg']
    bg.line.fill.background()
    
    # 装饰圆形
    add_decorative_elements(slide, prs, [MBE_COLORS['yellow'], MBE_COLORS['purple'], MBE_COLORS['red']])
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "大学新生开学第一课"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.font.color.rgb = MBE_COLORS['text']
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "开启精彩的大学之旅"
    p.font.size = Pt(24)
    p.font.name = 'Microsoft YaHei'
    p.font.color.rgb = MBE_COLORS['gray']
    
    # 右侧装饰 - 卡通人物区域（用形状代替）
    # 人物头部
    head = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(10), Inches(2.5),
        Inches(1.5), Inches(1.5)
    )
    head.fill.solid()
    head.fill.fore_color.rgb = MBE_COLORS['yellow']  # 黄色衣服 = 头部（简化表示）
    head.line.color.rgb = MBE_COLORS['border']
    head.line.width = Pt(3)
    
    # 人物身体
    body = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.5), Inches(4),
        Inches(2.5), Inches(2)
    )
    body.fill.solid()
    body.fill.fore_color.rgb = MBE_COLORS['blue']
    body.line.color.rgb = MBE_COLORS['border']
    body.line.width = Pt(3)
    
    # 对话气泡
    bubble = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.5), Inches(1.8),
        Inches(2), Inches(0.8)
    )
    bubble.fill.solid()
    bubble.fill.fore_color.rgb = MBE_COLORS['bg']
    bubble.line.color.rgb = MBE_COLORS['border']
    bubble.line.width = Pt(2)
    
    bubble_text = slide.shapes.add_textbox(Inches(8.7), Inches(1.9), Inches(1.6), Inches(0.6))
    tf = bubble_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Hello!"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = MBE_COLORS['text']
    
    # 底部信息
    bottom_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(5), Inches(0.5))
    tf = bottom_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2026年秋季学期"
    p.font.size = Pt(14)
    p.font.color.rgb = MBE_COLORS['gray']
    
    # ===== 第2页：目录 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = MBE_COLORS['bg']
    bg.line.fill.background()
    
    # 装饰
    add_decorative_elements(slide, prs, [MBE_COLORS['green'], MBE_COLORS['orange'], MBE_COLORS['purple']])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(5), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "今日议程"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.font.color.rgb = MBE_COLORS['text']
    
    # 目录项
    toc_items = [
        ("01", "认识校园", MBE_COLORS['yellow']),
        ("02", "学业规划", MBE_COLORS['purple']),
        ("03", "校园生活", MBE_COLORS['red']),
        ("04", "人际关系", MBE_COLORS['blue']),
        ("05", "未来展望", MBE_COLORS['green']),
    ]
    
    for i, (num, text, color) in enumerate(toc_items):
        y = 1.8 + i * 1
        
        # 序号圆圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8), Inches(y),
            Inches(0.7), Inches(0.7)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        
        # 序号文字
        num_box = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.1), Inches(0.7), Inches(0.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = MBE_COLORS['bg']
        p.alignment = PP_ALIGN.CENTER
        
        # 文字
        text_box = slide.shapes.add_textbox(Inches(1.8), Inches(y + 0.1), Inches(5), Inches(0.5))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(20)
        p.font.name = 'Microsoft YaHei'
        p.font.color.rgb = MBE_COLORS['text']
    
    # 右侧装饰插画区
    illustration_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.5), Inches(1.5),
        Inches(4), Inches(5)
    )
    illustration_box.fill.solid()
    illustration_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    illustration_box.line.color.rgb = MBE_COLORS['border']
    illustration_box.line.width = Pt(2)
    
    # ===== 第3页：认识校园 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = MBE_COLORS['bg']
    bg.line.fill.background()
    
    add_decorative_elements(slide, prs, [MBE_COLORS['purple'], MBE_COLORS['red'], MBE_COLORS['blue']])
    
    # 左侧内容区
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(6), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "认识校园"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.font.color.rgb = MBE_COLORS['text']
    
    # 内容要点
    points = [
        ("教学楼分布", "了解各学院教学楼位置"),
        ("图书馆资源", "海量藏书与电子资源"),
        ("食堂指南", "各校区美食推荐"),
        ("校园地图", "快速熟悉校园环境"),
    ]
    
    colors = [MBE_COLORS['yellow'], MBE_COLORS['purple'], MBE_COLORS['red'], MBE_COLORS['blue']]
    for i, (title, desc) in enumerate(points):
        y = 1.8 + i * 1.2
        
        # 彩色圆点
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8), Inches(y + 0.1),
            Inches(0.4), Inches(0.4)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = colors[i]
        dot.line.fill.background()
        
        # 标题
        t_box = slide.shapes.add_textbox(Inches(1.4), Inches(y), Inches(5), Inches(0.5))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.name = 'Microsoft YaHei'
        p.font.color.rgb = MBE_COLORS['text']
        
        # 描述
        d_box = slide.shapes.add_textbox(Inches(1.4), Inches(y + 0.4), Inches(5), Inches(0.5))
        tf = d_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.name = 'Microsoft YaHei'
        p.font.color.rgb = MBE_COLORS['gray']
    
    # 右侧插画
    ill_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8), Inches(1.2),
        Inches(4.8), Inches(5.5)
    )
    ill_box.fill.solid()
    ill_box.fill.fore_color.rgb = RGBColor(250, 250, 250)
    ill_box.line.color.rgb = MBE_COLORS['border']
    ill_box.line.width = Pt(2)
    
    # 插画内容 - 校园建筑简化图
    building = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.5), Inches(3),
        Inches(2), Inches(2.5)
    )
    building.fill.solid()
    building.fill.fore_color.rgb = MBE_COLORS['blue']
    building.line.color.rgb = MBE_COLORS['border']
    building.line.width = Pt(2)
    
    # 旗杆
    flag = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(10.9), Inches(2.2),
        Inches(0.1), Inches(0.8)
    )
    flag.fill.solid()
    flag.fill.fore_color.rgb = MBE_COLORS['border']
    flag.line.fill.background()
    
    # ===== 第4页：学业规划 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = MBE_COLORS['bg']
    bg.line.fill.background()
    
    add_decorative_elements(slide, prs, [MBE_COLORS['orange'], MBE_COLORS['green'], MBE_COLORS['yellow']])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(6), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "学业规划"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.font.color.rgb = MBE_COLORS['text']
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    items = [
        "• 制定学习目标与学期计划",
        "• 合理安排课程与时间",
        "• 善用图书馆与自习室",
        "• 积极参加学术讲座",
        "• 培养自主学习能力",
    ]
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.name = 'Microsoft YaHei'
        p.font.color.rgb = MBE_COLORS['text']
        p.space_after = Pt(20)
    
    # 右侧
    ill_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8), Inches(1.2),
        Inches(4.8), Inches(5.5)
    )
    ill_box.fill.solid()
    ill_box.fill.fore_color.rgb = RGBColor(250, 250, 250)
    ill_box.line.color.rgb = MBE_COLORS['border']
    ill_box.line.width = Pt(2)
    
    # ===== 第5页：总结 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = MBE_COLORS['purple']
    bg.line.fill.background()
    
    # 装饰
    for i, (x, y, color) in enumerate([
        (1, 1, MBE_COLORS['yellow']),
        (11, 0.5, MBE_COLORS['red']),
        (10, 5.5, MBE_COLORS['blue']),
    ]):
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(2), Inches(2))
        c.fill.solid()
        c.fill.fore_color.rgb = color
        c.fill.fore_color.brightness = 0.3
        c.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0), Inches(2.5), Inches(13.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "祝大家大学生活愉快！"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.font.color.rgb = MBE_COLORS['bg']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0), Inches(4.2), Inches(13.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "开启属于你的精彩篇章"
    p.font.size = Pt(24)
    p.font.name = 'Microsoft YaHei'
    p.font.color.rgb = RGBColor(220, 200, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 保存
    if output_path is None:
        output_path = os.path.join(os.path.expanduser("~"), "Desktop", "大学新生开学第一课.pptx")
    
    prs.save(output_path)
    print(f"✅ MBE风格PPT已生成: {output_path}")
    return output_path

if __name__ == "__main__":
    output = create_mbe_ppt()
    print(f"\n📁 文件位置: {output}")
