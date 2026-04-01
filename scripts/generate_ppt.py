#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
PPT Generator - 智能 PPT 生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.dml.color import RGBColor as RgbColor
import os

# 配色方案
COLOR_PALETTES = {
    'midnight-executive': {
        'name': '午夜商务',
        'primary': (30, 39, 97),      # #1E2761 深蓝
        'secondary': (202, 220, 252), # #CADCFC 冰蓝
        'accent': (255, 255, 255),    # #FFFFFF 白
        'bg_light': (255, 255, 255),  # 白色背景
        'text': (30, 39, 97),         # 深蓝文字
        'text_light': (107, 114, 128) # 灰色文字
    },
    'tech-dark': {
        'name': '科技深空',
        'primary': (13, 17, 23),      # #0D1117 深黑
        'secondary': (22, 27, 34),     # #161B22 深灰
        'accent': (88, 166, 255),      # #58A6FF 科技蓝
        'bg_light': (255, 255, 255),
        'text': (240, 246, 252),      # 白色文字
        'text_light': (139, 148, 158) # 灰色文字
    },
    'coral-energy': {
        'name': '珊瑚活力',
        'primary': (249, 97, 103),    # #F96167 珊瑚红
        'secondary': (249, 231, 149), # #F9E795 金黄
        'accent': (47, 60, 126),      # #2F3C7E 深海蓝
        'bg_light': (255, 255, 255),
        'text': (31, 41, 55),         # 深灰文字
        'text_light': (107, 114, 128)
    }
}

def create_presentation(title="演示文稿", subtitle="AI 智能生成", palette='midnight-executive', output_path=None):
    """创建演示文稿"""
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 宽屏
    prs.slide_height = Inches(7.5)
    
    colors = COLOR_PALETTES.get(palette, COLOR_PALETTES['midnight-executive'])
    
    # 1. 封面页
    create_cover_slide(prs, title, subtitle, colors)
    
    # 2. 目录页
    create_toc_slide(prs, colors)
    
    # 3. 内容页 1 - 概览
    create_content_slide(prs, "内容概览", [
        "第一个核心要点，包含详细说明内容",
        "第二个核心要点，展示关键信息",
        "第三个核心要点，补充相关细节"
    ], colors)
    
    # 4. 内容页 2 - 大数据展示
    create_big_number_slide(prs, "关键数据", [
        ("85%", "用户增长"),
        ("3.2倍", "效率提升"),
        ("200+", "客户数量")
    ], colors)
    
    # 5. 内容页 3 - 要点列表
    create_content_slide(prs, "核心要点", [
        "要点一：突出重点内容和关键信息",
        "要点二：展示核心优势和独特价值",
        "要点三：阐述主要功能和解决方案",
        "要点四：强调实际效果和成果"
    ], colors)
    
    # 6. 内容页 4 - 案例分析
    create_content_slide(prs, "案例分析", [
        "案例一：某科技公司通过智能系统提升效率 40%",
        "案例二：某电商平台月活跃用户突破 100 万",
        "案例三：某金融机构风险控制准确率提升至 98%"
    ], colors)
    
    # 7. 内容页 5 - 总结
    create_summary_slide(prs, [
        "核心价值：降本增效，提升竞争力",
        "关键优势：技术领先，服务专业",
        "未来展望：持续创新，合作共赢"
    ], colors)
    
    # 8. 结束页
    create_end_slide(prs, colors)
    
    # 保存文件
    if output_path is None:
        output_path = os.path.join(os.path.expanduser("~"), "Desktop", f"{title}.pptx")
    
    prs.save(output_path)
    print(f"✅ PPT 已生成: {output_path}")
    return output_path

def create_cover_slide(prs, title, subtitle, colors):
    """创建封面页"""
    blank_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(blank_layout)
    
    # 深色背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(*colors['primary'])
    background.line.fill.background()
    
    # 装饰圆形 - 右上角
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5), Inches(5), Inches(5)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = RgbColor(*colors['secondary'])
    circle1.line.fill.background()
    
    # 装饰圆形 - 左下角
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(-1), Inches(5), Inches(4), Inches(4)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RgbColor(*colors['secondary'])
    circle2.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.5), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RgbColor(*colors['accent'])
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.8))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = RgbColor(*colors['secondary'])
        p.alignment = PP_ALIGN.CENTER
    
    # 底部装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(5.8), Inches(2.5), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RgbColor(*colors['accent'])
    line.line.fill.background()

def create_toc_slide(prs, colors):
    """创建目录页"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # 白色背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(*colors['bg_light'])
    background.line.fill.background()
    
    # 页面标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "目录"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RgbColor(*colors['primary'])
    
    # 左侧装饰条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(0.1), Inches(5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RgbColor(*colors['primary'])
    bar.line.fill.background()
    
    # 目录项
    toc_items = ["内容概览", "关键数据", "核心要点", "案例分析", "总结与展望"]
    for i, item in enumerate(toc_items):
        # 序号
        num_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.8 + i * 0.9), Inches(0.6), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"0{i+1}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RgbColor(*colors['primary'])
        
        # 标题
        text_box = slide.shapes.add_textbox(Inches(2.0), Inches(1.8 + i * 0.9), Inches(9), Inches(0.6))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RgbColor(*colors['text'])
        
        # 分隔线
        if i < len(toc_items) - 1:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(2.0), Inches(2.5 + i * 0.9), Inches(9), Inches(0.01)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RgbColor(229, 229, 229)
            line.line.fill.background()

def create_content_slide(prs, title, content_list, colors):
    """创建内容页"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # 白色背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(*colors['bg_light'])
    background.line.fill.background()
    
    # 顶部装饰条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = RgbColor(*colors['primary'])
    top_bar.line.fill.background()
    
    # 页面标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(*colors['primary'])
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = RgbColor(*colors['text'])
        p.space_after = Pt(16)

def create_big_number_slide(prs, title, numbers, colors):
    """创建大数据展示页"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # 白色背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(*colors['bg_light'])
    background.line.fill.background()
    
    # 顶部装饰条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = RgbColor(*colors['primary'])
    top_bar.line.fill.background()
    
    # 页面标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(*colors['primary'])
    
    # 数据卡片
    card_width = Inches(3.5)
    card_height = Inches(2.8)
    start_x = Inches(0.8)
    start_y = Inches(2.2)
    gap = Inches(0.5)
    
    for i, (value, label) in enumerate(numbers):
        x = start_x + i * (card_width + gap)
        
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RgbColor(*colors['primary'])
        card.line.fill.background()
        
        # 数字
        value_box = slide.shapes.add_textbox(x, start_y + Inches(0.6), card_width, Inches(1.2))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RgbColor(*colors['accent'])
        p.alignment = PP_ALIGN.CENTER
        
        # 标签
        label_box = slide.shapes.add_textbox(x, start_y + Inches(1.8), card_width, Inches(0.6))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = RgbColor(*colors['secondary'])
        p.alignment = PP_ALIGN.CENTER

def create_summary_slide(prs, points, colors):
    """创建总结页"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # 深色背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(*colors['primary'])
    background.line.fill.background()
    
    # 页面标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "总结"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RgbColor(*colors['accent'])
    
    # 要点
    for i, point in enumerate(points):
        # 序号圆点
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.8), Inches(1.8 + i * 1.2 + 0.1), Inches(0.4), Inches(0.4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RgbColor(*colors['accent'])
        circle.line.fill.background()
        
        # 序号
        num_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8 + i * 1.2 + 0.05), Inches(0.4), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RgbColor(*colors['primary'])
        p.alignment = PP_ALIGN.CENTER
        
        # 文字
        text_box = slide.shapes.add_textbox(Inches(1.4), Inches(1.8 + i * 1.2), Inches(10.5), Inches(0.8))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = RgbColor(*colors['accent'])

def create_end_slide(prs, colors):
    """创建结束页"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # 深色背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(*colors['primary'])
    background.line.fill.background()
    
    # 装饰圆形
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(4.5), Inches(2), Inches(4.5), Inches(4.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RgbColor(*colors['secondary'])
    circle.line.fill.background()
    
    # 结束语
    end_box = slide.shapes.add_textbox(Inches(0), Inches(3.2), prs.slide_width, Inches(1.5))
    tf = end_box.text_frame
    p = tf.paragraphs[0]
    p.text = "谢谢观看"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RgbColor(*colors['accent'])
    p.alignment = PP_ALIGN.CENTER


if __name__ == "__main__":
    # 生成示例 PPT
    output = create_presentation(
        title="AI发展趋势报告",
        subtitle="2026年度行业分析",
        palette='midnight-executive',
        output_path=None
    )
    print(f"\n📁 文件位置: {output}")
